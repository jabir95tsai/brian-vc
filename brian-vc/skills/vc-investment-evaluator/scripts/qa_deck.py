#!/usr/bin/env python3
"""Mode-aware structural PPTX QA；程式化檢查不取代視覺 QA。"""

from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path

if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8")
    sys.stderr.reconfigure(encoding="utf-8")

try:
    from pptx import Presentation
except ModuleNotFoundError as exc:
    if exc.name != "pptx":
        raise
    print("ERROR: qa_deck.py requires python-pptx; install brian-vc/requirements.txt", file=sys.stderr)
    raise SystemExit(2)


def slide_text(slide) -> str:
    return "\n".join(shape.text_frame.text for shape in slide.shapes if shape.has_text_frame)


def slide_tables(slide):
    return [shape.table for shape in slide.shapes if getattr(shape, "has_table", False)]


def find_slides(prs, keywords):
    return [(index, slide) for index, slide in enumerate(prs.slides) if any(key in slide_text(slide) for key in keywords)]


def data_rows(table) -> int:
    return max(0, len(table.rows) - 1)


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("pptx", type=Path)
    parser.add_argument("--full", action="store_true", help="apply Full-critical additions")
    parser.add_argument("--mode", choices=("full", "degraded", "blocked", "quick-screen"), default="full")
    parser.add_argument("--min-slides", type=int)
    parser.add_argument("--json-output", type=Path)
    args = parser.parse_args()

    prs = Presentation(args.pptx)
    checks: list[dict] = []

    def record(name: str, ok: bool, detail: str = "") -> bool:
        checks.append({"name": name, "status": "pass" if ok else "fail", "detail": detail})
        return ok

    thresholds = {
        "full": {"executive": 15, "full": 18, "comps": 5, "team": 2, "history": 8, "quant": 2},
        "degraded": {"executive": 10, "full": 15, "comps": 1, "team": 1, "history": 5, "quant": 1},
        "blocked": {"executive": 10, "full": 15, "comps": 1, "team": 1, "history": 5, "quant": 1},
        "quick-screen": {"executive": 3, "full": 3, "comps": 0, "team": 0, "history": 0, "quant": 0},
    }[args.mode]
    minimum = args.min_slides or thresholds["full" if args.full else "executive"]
    record("檔案可開啟且頁數足夠", len(prs.slides) >= minimum, f"{len(prs.slides)} 頁（門檻 {minimum}）")

    def table_page(label, keywords, min_rows=1, min_tables=1, gap_keywords=()):
        hits = find_slides(prs, keywords)
        if not hits:
            return record(label, False, f"找不到含 {keywords} 的頁")
        tables = max((slide_tables(slide) for _, slide in hits), key=len)
        rows = max((data_rows(table) for table in tables), default=0)
        text = "\n".join(slide_text(slide) for _, slide in hits)
        gap_ok = args.mode != "full" and gap_keywords and any(key in text for key in gap_keywords)
        ok = (len(tables) >= min_tables and rows >= min_rows) or bool(gap_ok)
        return record(label, ok, f"表格 {len(tables)}；資料列 {rows}；degraded 缺件揭露={bool(gap_ok)}")

    if args.mode != "quick-screen":
        table_page("同業比較／缺件揭露", ["同業比較", "同業估值", "可比公司", "已驗證同業"], thresholds["comps"], gap_keywords=("尚待官方", "尚未完成", "待補"))
        table_page("IRR+Multiple 雙矩陣", ["報酬率", "IRR", "Return Multiple"], 3, 2, gap_keywords=("尚未取得本輪條件", "本輪條件未揭露", "blocked", "敏感度"))
        table_page("財務預測逐項多年", ["財務預測", "財測", "獨立基準", "基準情境"], 5, gap_keywords=("無可用財測假設", "獨立財測未建立", "blocked"))
        table_page("經營團隊／缺件揭露", ["經營團隊", "團隊", "治理資料"], thresholds["team"], gap_keywords=("未提供", "尚缺", "待補"))
        table_page("歷史財務逐項", ["損益", "歷史 IS", "營業收入", "歷史財務"], thresholds["history"])
        table_page("技術頁有量化表", ["競品", "技術", "供應韌性"], thresholds["quant"])
        table_page("產業／市場頁有量化表", ["市場規模", "產業", "市場證據"], thresholds["quant"])
    else:
        record("初篩必含結論邊界", bool(find_slides(prs, ["初篩", "資料不足", "下一步", "不代表投資決策"])))

    if args.full and args.mode != "quick-screen":
        table_page("風險矩陣頁", ["風險矩陣", "風險分析", "風險"], 3 if args.mode == "full" else 1)
        record("RedTeam／反方頁存在", bool(find_slides(prs, ["RedTeam", "紅隊", "反對", "靈魂拷問", "反方"])))
        record("資料缺失／補件頁存在", bool(find_slides(prs, ["資料缺失", "補件", "缺件", "資料不足", "前置條件"])))

    report = {
        "status": "pass" if all(item["status"] == "pass" for item in checks) else "fail",
        "mode": args.mode,
        "variant": "full-critical" if args.full else "executive",
        "pptx": str(args.pptx.resolve()),
        "slide_count": len(prs.slides),
        "checks": checks,
        "visual_qa_required": True,
    }
    if args.json_output:
        args.json_output.parent.mkdir(parents=True, exist_ok=True)
        args.json_output.write_text(json.dumps(report, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    print(json.dumps(report, ensure_ascii=False, indent=2))
    return 0 if report["status"] == "pass" else 1


if __name__ == "__main__":
    raise SystemExit(main())
